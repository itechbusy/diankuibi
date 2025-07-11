import asyncio
import base64
import json
import logging
import os
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from typing import List

import aiofiles
import fitz
from docx import Document
from openai import AsyncOpenAI
from openai import OpenAI
from openpyxl.reader.excel import load_workbook
from pptx import Presentation
from pydantic import UUID4

from common.str_transcoding import str_decrypt
from processor.models.image_model import MultiplePictureModel
from processor.models.model_settings import get_model, get_default_model
from processor.prompt_templates import *
from task_flow.models import ImageInfo

logging = logging.getLogger('processor')


async def chat(question=None, model_name=None):
    return invoke_text_model(question, model_name, TEXT_CHAT_SYSTEM)


async def text_reasoning(prompt=None, model_name=None):
    return await invoke_text_model(prompt, model_name, TEXT_REASON_SYSTEM)


async def invoke_text_model(prompt, model_name, sys_template):
    message = [{"role": "system", "content": f"{sys_template}"}, {"role": "user", "content": f"{prompt}"}]

    model_info = await get_default_model(0) if model_name is None else await get_model(model_name=model_name)
    if model_info is None:
        logging.error(f"No model fits the criteria. model type 0 , model name {model_name}")
        raise ValueError(f"No model fits the criteria. model type 0 , model name {model_name}")
    return await invoke_model(model_info, json.dumps(message))


async def picture_reasoning(data: MultiplePictureModel, prompt=None, model_name=None):
    if not data: raise ValueError("Parameter 'data' is empty.")
    if not prompt: raise ValueError("Prompt 'data' is empty.")
    content = []
    if data.url:
        content.append({"type": "image_url", "image_url": {"url": data.url}})

    else:
        if data.base64_data:
            if not data.image_type:
                raise ValueError(f"Parameter 'image_type' is empty. Image id {data.image_id} ")
        content.append(
            {"type": "image_url", "image_url": {"url": f'data:image/{data.image_type};base64,{data.base64_data}'}})

    content.append({"type": "text", "text": f"{prompt}"})

    model_info = await get_default_model(1) if model_name is None else await get_model(model_name=model_name)
    if model_info is None:
        logging.error(f"No model fits the criteria. model type 1 , model name {model_name}")
        raise ValueError(f"No model fits the criteria. model type 1 , model name {model_name}")

    return invoke_model(model_info, json.dumps([{"role": "user", "content": content}]))

    # dict_data = asyncio.run(json_response_to_dict(response))
    # if type(dict_data) is not dict:
    #     return PictureReasoningResult(image_id=data.image_id, successfully=False)
    # else:
    #     return PictureReasoningResult(image_id=data.image_id, interpretation=dict_data["r"], classify=dict_data["s"])


def extract_json_content(text):
    start_index = text.find("```json")
    if start_index == -1:
        return None
    end_index = text.find("```", start_index + len("```json"))
    if end_index == -1:
        return None
    return text[start_index:end_index + len("```")]


def json_response_to_dict(text):
    cleaned_str = text.strip('```json').strip()
    try:
        return json.loads(cleaned_str)
    except json.JSONDecodeError as e:
        logging.error(f"Picture reasoning fail. text :{text}, error :{e}")
        return None


async def multiple_picture_reasoning(data: List[MultiplePictureModel], prompt=None, model_name=None):
    if not data: raise ValueError("Parameter 'data' is empty.")
    if not prompt: raise ValueError("Prompt 'data' is empty.")
    results = []
    with ThreadPoolExecutor(max_workers=5) as executor:
        futures = {executor.submit(picture_reasoning, task, prompt, model_name): task for task in data}
        for future in as_completed(futures):
            try:
                results.append(future.result())
            except Exception as e:
                logging.error(f"Picture reasoning fail., error :{e}")

    return results


async def invoke_model(client_info, message):
    client = AsyncOpenAI(
        api_key=client_info.api_key,
        base_url=client_info.base_url,
        timeout=client_info.timeout,
        max_retries=client_info.max_retries,
    )

    completion = await client.chat.completions.create(
        model=client_info.model_name,
        temperature=client_info.temperature,
        messages=json.loads(message),
        extra_body={"enable_thinking": False}
    )
    return await extract_conversation_content(completion)


async def extract_conversation_content(answer):
    text = answer.choices[0].message.content
    if answer.choices[0].finish_reason == "stop":
        print(text)
        return text
    else:
        logging.warn(f'Incomplete answer, answer {text}')
        return None


async def generate_image_description(image_path, picture_reasoning_prompt, picture_reasoning_model_id):
    """Call Qwen VL to generate image descriptions (supports local files)"""
    try:
        # Base64 encoding processing (needs to adapt to different image formats)
        with aiofiles.open(image_path, "rb") as img_file:
            base64_image = base64.b64encode(await img_file.read()).decode('utf-8')
            mime_type = "image/jpeg" if image_path.lower().endswith('.jpg') else "image/png"

        response = picture_reasoning(MultiplePictureModel(
            image_id=UUID4.hex,
            base64_data=base64_image,
            image_type=mime_type
        ), prompt=str_decrypt(picture_reasoning_prompt), model_name=picture_reasoning_model_id)

        return f"<--AI图片描述 AI image description：{response}-->"
    except Exception as e:
        logging.error(f"Call model failed : {str(e)}")
        return ""


def extract_and_process_images(file_path, output_dir, picture_reasoning_prompt, picture_reasoning_model_id):
    """Extract embedded images, record context, generate descriptions and record them in the database"""
    file_ext = os.path.splitext(file_path)[1].lower()
    images = []
    context_records = []

    try:
        if file_ext in ['.xlsx', '.xls']:
            wb = load_workbook(file_path)
            for sheet_idx, sheet in enumerate(wb.worksheets):
                # Collect only the first row of the worksheet as context
                sheet_context = []
                for row in sheet.iter_rows(max_row=1):
                    row_data = [str(cell.value or "") for cell in row]
                    sheet_context.append(f"Row1: {' | '.join(row_data)}")
                context_text = "\n".join(sheet_context)[:300]  # 限制长度

                for img_idx, image in enumerate(sheet._images):
                    # Save Picture
                    img_bytes = image._data()
                    img_name = f"excel_{sheet.title}_img{img_idx}.png"
                    img_path = os.path.join(output_dir, img_name)
                    with open(img_path, "wb") as f:
                        f.write(img_bytes)
                    images.append(img_path)

                    # Record context
                    context_records.append({
                        'image_path': img_path,
                        'context_type': 'excel_sheet',
                        'context_data': {
                            'sheet_name': sheet.title,
                            'preview': context_text,
                            'position': f"单元格范围: {image.anchor._from}"
                        }
                    })

        elif file_ext in ['.docx', '.doc']:
            doc = Document(file_path)
            # Establish a mapping between paragraphs and images
            para_image_map = {}
            for i, paragraph in enumerate(doc.paragraphs):
                for run in paragraph.runs:
                    if run.element.xpath('.//w:drawing'):
                        for rel_id in run.element.xpath('.//a:blip/@r:embed'):
                            rel = doc.part.rels.get(rel_id)
                            if rel and "image" in rel.reltype:
                                parent_paragraph = paragraph
                                para_text = ' '.join([run.text for run in parent_paragraph.runs])
                                ext = rel.target_part.content_type.split('/')[-1]
                                img_name = f"word_para_img{len(images)}.{ext}"
                                img_path = os.path.join(output_dir, img_name)
                                with open(img_path, "wb") as f:
                                    f.write(rel.target_part.blob)
                                images.append(img_path)
                                context_records.append({
                                    'image_path': img_path,
                                    'context_type': 'word_paragraph',
                                    'context_data': {
                                        'text': para_text[:500],
                                        'style': parent_paragraph.style.name,
                                        'position': f"段落位置: {i + 1}"
                                    }
                                })

        elif file_ext in ['.pptx', '.ppt']:
            prs = Presentation(file_path)
            for slide_idx, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text = shape.text.splitlines()[:1]
                        break
                context_text = "\n".join(slide_text)[:300]

                for shape_idx, shape in enumerate(slide.shapes):
                    if hasattr(shape, "image"):
                        img_data = shape.image.blob
                        img_name = f"ppt_slide{slide_idx}_img{shape_idx}.{shape.image.ext}"
                        img_path = os.path.join(output_dir, img_name)
                        with open(img_path, "wb") as f:
                            f.write(img_data)

                        images.append(img_path)
                        context_records.append({
                            'image_path': img_path,
                            'context_type': 'ppt_slide',
                            'context_data': {
                                'slide_number': slide_idx + 1,
                                'content': context_text,
                                'notes': slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
                            }
                        })
        elif file_ext in ['.pdf']:
            doc = fitz.open(file_path)
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                page_text = page.get_text()[:300]
                image_list = page.get_images(full=True)
                for img_idx, img in enumerate(image_list, start=1):
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    img_ext = base_image["ext"]
                    img_name = f"pdf_page{page_num + 1}_img{img_idx}.{img_ext}"
                    img_path = os.path.join(output_dir, img_name)
                    with open(img_path, "wb") as f:
                        f.write(image_bytes)
                    images.append(img_path)

                    context_records.append({
                        'image_path': img_path,
                        'context_type': 'pdf_page',
                        'context_data': {
                            'page_number': page_num + 1,
                            'content': page_text
                        }
                    })
            doc.close()

        # Generate image descriptions and record them in the database
        document_name = os.path.basename(file_path)
        for record in context_records:
            image_path = record['image_path']
            context_text = str(record['context_data'])
            description = asyncio.run(
                generate_image_description(image_path, picture_reasoning_prompt, picture_reasoning_model_id))
            ImageInfo.objects.create(
                document_name=document_name,
                image_path=image_path,
                context_text=context_text,
                image_description=description
            )

        return images, context_records

    except Exception as e:
        # Clean up created image files
        for img in images:
            if os.path.exists(img):
                os.remove(img)
        raise


def document_understanding(file_path, user_question, model_name=None):
    client_info = asyncio.run(get_default_model(0)) if model_name is None else asyncio.run(
        get_model(model_name=model_name))
    if client_info is None:
        logging.error(f"No model fits the criteria. model type 0 , model name {model_name}")
        raise ValueError(f"No model fits the criteria. model type 0 , model name {model_name}")
    base_model_name = client_info.model_name
    client = OpenAI(
        api_key=client_info.api_key,
        base_url=client_info.base_url,
        timeout=client_info.timeout,
        max_retries=client_info.max_retries,
    )

    messages = [
        {'role': 'system', 'content': 'You are a helpful assistant.'}
    ]

    # Select the processing method based on the model name
    if base_model_name == "qwen-long":
        # Qwen Long model uses file upload method
        file_object = client.files.create(file=Path(file_path), purpose="file-extract")
        file_id = file_object.id
        if file_id:
            messages.append({'role': 'system', 'content': f'fileid://{file_id}'})
    else:
        # Other models read the file content and use it as prompt words
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                file_content = file.read()
            messages.append({'role': 'system', 'content': f'文件内容：\n{file_content}'})
        except Exception as e:
            logging.error(f"fail to read file: {str(e)}")
            raise ValueError(f"读取文件失败 fail to read file: {str(e)}")

    messages.append({'role': 'user', 'content': user_question})

    completion = client.chat.completions.create(
        model=client_info.model_name,
        temperature=client_info.temperature,
        messages=messages,
        stream=True,
        stream_options={"include_usage": True}
    )

    full_content = ""
    for chunk in completion:
        if chunk.choices and chunk.choices[0].delta.content:
            full_content += chunk.choices[0].delta.content
            print(chunk.model_dump())

    return full_content


def document_understanding_text(user_question, model_name):
    try:
        client_info = asyncio.run(get_default_model(0)) if model_name is None else asyncio.run(
            get_model(model_name=model_name))
    except Exception as e:
        logging.error(f"Error occurred while getting model info: {e}")
        raise ValueError(f"Error getting model info: {e}")

    if client_info is None:
        logging.error(f"No model fits the criteria. model type 0 , model name {model_name}")
        raise ValueError(f"No model fits the criteria. model type 0 , model name {model_name}")

    try:
        client = OpenAI(
            api_key=client_info.api_key,
            base_url=client_info.base_url,
            timeout=client_info.timeout,
            max_retries=client_info.max_retries,
        )
    except Exception as e:
        logging.error(f"Error creating OpenAI client: {e}")
        raise ValueError(f"Error creating OpenAI client: {e}")

    try:
        completion = client.chat.completions.create(
            model=model_name,
            messages=[
                {'role': 'system', 'content': TEXT_REASON_SYSTEM},
                {'role': 'user', 'content': user_question}
            ],
            stream=True
        )
        # Process streaming response
        response_text = ""
        for chunk in completion:
            if chunk.choices[0].delta.content:
                response_text += chunk.choices[0].delta.content
        return response_text
    except Exception as e:
        logging.error(f"Error making OpenAI API call: {e}")
        raise ValueError(f"Error making OpenAI API call: {e}")
